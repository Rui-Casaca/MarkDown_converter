"""Shared pytest fixtures for doc2md tests."""

from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def docx_file(tmp_path: Path) -> Path:
    """Create a small but representative DOCX document."""
    from docx import Document

    document = Document()
    document.add_heading("Title One", level=1)
    document.add_paragraph("Hello world paragraph.")
    document.add_paragraph("Bullet item", style="List Bullet")

    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "H1"
    table.rows[0].cells[1].text = "H2"
    table.rows[1].cells[0].text = "v1"
    table.rows[1].cells[1].text = "v2"

    path = tmp_path / "sample.docx"
    document.save(str(path))
    return path


@pytest.fixture
def pptx_file(tmp_path: Path) -> Path:
    """Create a small PPTX presentation with a titled slide and body text."""
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Body line"

    path = tmp_path / "sample.pptx"
    presentation.save(str(path))
    return path


@pytest.fixture
def rich_docx_file(tmp_path: Path) -> Path:
    """Create a DOCX with bold, italic, and a hyperlink in one paragraph."""
    import docx.oxml.ns as ns
    from docx import Document
    from docx.oxml import OxmlElement

    document = Document()
    paragraph = document.add_paragraph()
    paragraph.add_run("normal ")
    bold_run = paragraph.add_run("bolded")
    bold_run.bold = True
    paragraph.add_run(" and ")
    italic_run = paragraph.add_run("slanted")
    italic_run.italic = True
    paragraph.add_run(" ")

    rel_id = document.part.relate_to(
        "https://example.com",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(ns.qn("r:id"), rel_id)
    new_run = OxmlElement("w:r")
    new_run.append(OxmlElement("w:rPr"))
    text_element = OxmlElement("w:t")
    text_element.text = "the link"
    new_run.append(text_element)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)

    path = tmp_path / "rich.docx"
    document.save(str(path))
    return path


@pytest.fixture
def rich_pptx_file(tmp_path: Path) -> Path:
    """Create a PPTX whose body has bold, italic, and a hyperlinked run."""
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"

    paragraph = slide.placeholders[1].text_frame.paragraphs[0]
    bold_run = paragraph.add_run()
    bold_run.text = "Bolded"
    bold_run.font.bold = True
    paragraph.add_run().text = " middle "
    italic_run = paragraph.add_run()
    italic_run.text = "Slanted"
    italic_run.font.italic = True
    paragraph.add_run().text = " "
    link_run = paragraph.add_run()
    link_run.text = "linked"
    link_run.hyperlink.address = "https://example.org"

    path = tmp_path / "rich.pptx"
    presentation.save(str(path))
    return path


@pytest.fixture
def sample_png(tmp_path: Path) -> Path:
    """Create a tiny PNG image file."""
    from PIL import Image

    path = tmp_path / "red.png"
    Image.new("RGB", (8, 8), (200, 30, 30)).save(str(path))
    return path


@pytest.fixture
def docx_with_comments(tmp_path: Path) -> Path:
    """Create a DOCX with a single review comment anchored to a run of text."""
    from docx import Document

    document = Document()
    document.add_heading("Quarterly Report", level=1)
    paragraph = document.add_paragraph()
    paragraph.add_run("The quarterly revenue ")
    commented = paragraph.add_run("grew by 12%")
    paragraph.add_run(" in Q2.")

    if not hasattr(document, "add_comment"):
        pytest.skip("python-docx build lacks comment authoring support")

    document.add_comment(
        runs=[commented],
        text="Please verify this figure.",
        author="Jane Doe",
        initials="JD",
    )

    path = tmp_path / "commented.docx"
    document.save(str(path))
    return path


@pytest.fixture
def docx_with_image(tmp_path: Path, sample_png: Path) -> Path:
    """Create a DOCX containing a single embedded picture."""
    from docx import Document

    document = Document()
    document.add_paragraph("Before image")
    document.add_picture(str(sample_png))
    document.add_paragraph("After image")

    path = tmp_path / "with_image.docx"
    document.save(str(path))
    return path


@pytest.fixture
def pptx_with_image(tmp_path: Path, sample_png: Path) -> Path:
    """Create a PPTX containing a single embedded picture."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(sample_png), Inches(1), Inches(1))

    path = tmp_path / "with_image.pptx"
    presentation.save(str(path))
    return path
